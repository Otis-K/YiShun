from pathlib import Path
import os
import shutil
import zipfile

WORKSPACE_TEMP = Path(__file__).resolve().parents[2] / "work" / ".tmp"
WORKSPACE_TEMP.mkdir(parents=True, exist_ok=True)
os.environ["TEMP"] = str(WORKSPACE_TEMP)
os.environ["TMP"] = str(WORKSPACE_TEMP)

from docx import Document
from docx.shared import Inches
from openpyxl import Workbook
from openpyxl.drawing.image import Image as ExcelImage
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as PptInches

import matrix_runner as matrix


TARGET_BYTES = 10_000_000
MEDIA_COUNT = 20
FIXTURE_VERSION = "8-direct-ppt-png-boundary"
ROOT = matrix.ACCEPTANCE / "samples" / "L3" / "office-boundary"
IMAGE_ROOT = ROOT / "images-supported-10mb"
PPT_IMAGE_ROOT = ROOT / "images-supported-ppt-direct-5mb"


def make_images():
    IMAGE_ROOT.mkdir(parents=True, exist_ok=True)
    paths = []
    for index in range(MEDIA_COUNT):
        path = IMAGE_ROOT / f"image-{index:03d}.png"
        if not path.exists() or path.stat().st_size < 400_000:
            width = height = 384
            image = Image.frombytes("RGB", (width, height), os.urandom(width * height * 3))
            image.save(path, compress_level=0)
        paths.append(path)
    return paths


def make_docx(images, path):
    document = Document()
    document.add_heading("Tool Plus L3 Office Boundary", 0)
    document.add_paragraph("OLD Tool semantic marker")
    for index, image in enumerate(images):
        paragraph = document.add_paragraph(f"Embedded media {index + 1:03d} ")
        paragraph.add_run().add_picture(str(image), width=Inches(0.25))
    document.save(path)


def make_xlsx(images, path):
    workbook = Workbook()
    for index in range(100):
        sheet = workbook.active if index == 0 else workbook.create_sheet()
        sheet.title = f"Sheet{index + 1:03d}"
        sheet["A1"] = "OLD Tool semantic marker"
        sheet["A2"] = index + 1
        if index < len(images):
            picture = ExcelImage(str(images[index]))
            picture.width = picture.height = 24
            sheet.add_image(picture, "C2")
    workbook.save(path)


def make_pptx(images, path, slide_count=500):
    PPT_IMAGE_ROOT.mkdir(parents=True, exist_ok=True)
    ppt_images = []
    for index, _source in enumerate(images):
        resized = PPT_IMAGE_ROOT / f"image-{index:03d}.png"
        if not resized.exists() or resized.stat().st_size < 150_000:
            image = Image.frombytes("RGB", (256, 256), os.urandom(256 * 256 * 3))
            image.save(resized, compress_level=0)
        ppt_images.append(resized)
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for index in range(slide_count):
        slide = presentation.slides.add_slide(blank)
        text = slide.shapes.add_textbox(PptInches(0.2), PptInches(0.2), PptInches(4), PptInches(0.4))
        text.text_frame.text = f"OLD Tool slide {index + 1:03d}"
        if index < len(ppt_images):
            slide.shapes.add_picture(str(ppt_images[index]), PptInches(5), PptInches(0.2), width=PptInches(0.3), height=PptInches(0.3))
    presentation.save(path)


def append_padding(source, destination, content_bytes):
    shutil.copy2(source, destination)
    info = zipfile.ZipInfo("customXml/l3-boundary-padding.xml")
    info.compress_type = zipfile.ZIP_STORED
    prefix, suffix = b"<padding>", b"</padding>"
    if content_bytes < len(prefix) + len(suffix):
        raise AssertionError("padding content is too small")
    with zipfile.ZipFile(destination, "a", allowZip64=True) as archive:
        with archive.open(info, "w", force_zip64=True) as stream:
            stream.write(prefix)
            remaining = content_bytes - len(prefix) - len(suffix)
            block = b"A" * (1024 * 1024)
            while remaining:
                amount = min(remaining, len(block))
                stream.write(block[:amount])
                remaining -= amount
            stream.write(suffix)


def pad_exact(base, destination):
    trial = destination.with_suffix(destination.suffix + ".trial")
    append_padding(base, trial, 32)
    overhead = trial.stat().st_size - base.stat().st_size - 32
    trial.unlink()
    content_bytes = TARGET_BYTES - base.stat().st_size - overhead
    append_padding(base, destination, content_bytes)
    if destination.stat().st_size != TARGET_BYTES:
        raise AssertionError(f"{destination.name} size {destination.stat().st_size} != {TARGET_BYTES}")


def count_parts(path):
    with zipfile.ZipFile(path) as archive:
        names = [name.replace("\\", "/").lower() for name in archive.namelist()]
        bad = archive.testzip()
    if bad:
        raise AssertionError(f"corrupt ZIP member: {bad}")
    return {
        "sheets": sum(name.startswith("xl/worksheets/sheet") and name.endswith(".xml") for name in names),
        "slides": sum(name.startswith("ppt/slides/slide") and name.endswith(".xml") for name in names),
        "media": sum(name.startswith(("word/media/", "xl/media/", "ppt/media/")) for name in names),
    }


def main():
    ROOT.mkdir(parents=True, exist_ok=True)
    marker = ROOT / ".fixture-version"
    rebuild = not marker.exists() or marker.read_text(encoding="utf-8") != FIXTURE_VERSION
    smoke_docx = ROOT / "native-smoke.docx"
    smoke_xlsx = ROOT / "native-smoke.xlsx"
    smoke_pptx = ROOT / "native-smoke.pptx"
    if not smoke_docx.exists():
        document = Document()
        document.add_paragraph("Tool Plus native Office smoke test")
        document.save(smoke_docx)
    if not smoke_xlsx.exists():
        workbook = Workbook()
        workbook.active["A1"] = "Tool Plus native Office smoke test"
        workbook.save(smoke_xlsx)
    if not smoke_pptx.exists():
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(PptInches(0.2), PptInches(0.2), PptInches(4), PptInches(0.4)).text = "Tool Plus native Office smoke test"
        presentation.save(smoke_pptx)
    diagnostic_500 = ROOT / "ppt-500-slides-no-media.pptx"
    diagnostic_media = ROOT / "ppt-100-slides-20-media.pptx"
    if rebuild or not diagnostic_500.exists():
        make_pptx([], diagnostic_500, 500)
    if rebuild or not diagnostic_media.exists():
        make_pptx(make_images(), diagnostic_media, 100)
    images = make_images()
    builders = {"docx": make_docx, "xlsx": make_xlsx, "pptx": make_pptx}
    for extension, builder in builders.items():
        destination = ROOT / f"boundary-supported.{extension}"
        if rebuild or not destination.exists():
            builder(images, destination)
        limit = 5_000_000 if extension == "pptx" else TARGET_BYTES
        if destination.stat().st_size > limit:
            raise AssertionError(f"{destination.name} size {destination.stat().st_size} exceeds {limit}")
        print(destination, destination.stat().st_size, count_parts(destination), flush=True)
    marker.write_text(FIXTURE_VERSION, encoding="utf-8")


if __name__ == "__main__":
    main()
